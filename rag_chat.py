import streamlit as st
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains import RetrievalQA

from langchain_core.documents import Document
from langchain_text_splitters import CharacterTextSplitter

from docx import Document as DocxDocument
import fitz  # PyMuPDF
import os
import pandas as pd
import pptx
import json
import nest_asyncio

nest_asyncio.apply()
from dotenv import load_dotenv
load_dotenv()

GOOGLE_API_KEY = st.secrets.get("GOOGLE_API_KEY") or os.getenv("GOOGLE_API_KEY")

st.set_page_config(page_title="AskMyFile")

# --- App Title & Description ---
st.title("AskMyFile")
st.write("""
AskMyFile lets you upload a document and ask questions about it.  
Supported file types: PDF, DOCX, TXT, CSV, XLSX, PPTX, and JSON.
""")


@st.cache_resource
def load_model():
    """Load Gemini model with your Google API key"""
    return ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        temperature=0.5,
        top_p=0.9,
        max_output_tokens=2048,
        google_api_key=GOOGLE_API_KEY
    )


def format_dataframe_to_text(df):
    header = ", ".join(df.columns)
    rows = "\n".join([", ".join(map(str, row)) for row in df.values.tolist()])
    return f"{header}\n{rows}"


# --- Extract text from files ---
def extract_text_from_file(file):
    text = ""

    if file.type == "application/pdf":
        with fitz.open(stream=file.read(), filetype="pdf") as doc:
            for page in doc:
                text += page.get_text() + "\n"

    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = DocxDocument(file)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text.strip() + "\n"

    elif file.type == "text/plain":
        text = file.read().decode("utf-8")

    elif file.type == "text/csv":
        df = pd.read_csv(file)
        text = format_dataframe_to_text(df)

    elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        df = pd.read_excel(file)
        text = format_dataframe_to_text(df)

    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = pptx.Presentation(file)
        for slide_number, slide in enumerate(prs.slides):
            text += f"Slide {slide_number + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text.strip() + "\n"
            text += "\n"

    elif file.type == "application/json":
        content = file.read().decode("utf-8")
        try:
            json_obj = json.loads(content)
            text = json.dumps(json_obj, indent=2)
        except Exception as e:
            st.error(f"Failed to parse JSON: {e}")
            text = content  # fallback

    else:
        st.error("Unsupported file format.")

    return text.strip()


# --- Create Vectorstore ---
def create_vectorstore(text):
    if not text.strip():
        raise ValueError("No text found in the selected file to create vector store.")
    
    splitter = CharacterTextSplitter(chunk_size=300, chunk_overlap=50)
    docs = splitter.create_documents([text])
    
    if not docs:
        raise ValueError("Text splitting resulted in 0 chunks. Ensure file has meaningful content.")
    
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    return FAISS.from_documents(docs, embeddings)


# --- File Upload Section ---
file_types = {
    "PDF": "pdf",
    "Word Document (DOCX)": "docx",
    "Text File (TXT)": "txt",
    "CSV File": "csv",
    "Excel File (XLSX)": "xlsx",
    "PowerPoint (PPTX)": "pptx",
    "JSON File": "json"
}

selected_type = st.selectbox("Select the file type you're uploading", list(file_types.keys()))
selected_ext = file_types[selected_type]

uploaded_files = st.file_uploader(
    f"Upload one or more {selected_type} files",
    type=[selected_ext],
    accept_multiple_files=True
)

selected_file = None
file_texts = {}

if uploaded_files:
    st.subheader("📁 Select a file to preview and query:")
    filenames = [file.name for file in uploaded_files]
    selected_filename = st.selectbox("Choose a file", filenames)

    for file in uploaded_files:
        if file.name == selected_filename:
            with st.spinner(f"Extracting content from {file.name}..."):
                file_text = extract_text_from_file(file)
                file_texts[file.name] = file_text
                selected_file = file
            break

    if selected_file:
        if st.button("🔍 Preview File Content"):
            with st.expander("📄 File Preview", expanded=True):
                st.text_area("Extracted Content", file_texts[selected_file.name], height=300)

        # --- Create Vector Store and QA Chain ---
        vectorstore = create_vectorstore(file_texts[selected_file.name])
        retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
        llm = load_model()

        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            retriever=retriever,
            chain_type="stuff",
            return_source_documents=True
        )

        user_query = st.text_input("💬 Ask a question based on the selected file:")
        if user_query:
            with st.spinner("🔄 Processing your query..."):
                result = qa_chain.invoke({"query": user_query})

            st.markdown("### ✅ Answer")
            st.success(result["result"])

            with st.expander("🔎 Retrieved Chunks"):
                for doc in result["source_documents"]:
                    st.markdown(
                        f"<div style='background-color:#f1f8e9; padding:10px; border-radius:8px;'>{doc.page_content}</div>",
                        unsafe_allow_html=True
                    )
